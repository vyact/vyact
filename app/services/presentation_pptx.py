"""Editable, native PowerPoint renderer for AI-generated presentations."""
from __future__ import annotations

import base64
import re
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from services.presentation_design import prepare_presentation


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_FAMILY = "Noto Sans KR"
HEADER_TITLE_MAX_SIZE = 35
HEADER_TITLE_MEDIUM_SIZE = 30
HEADER_TITLE_LONG_SIZE = 25
HEADER_TITLE_MEDIUM_LENGTH = 34
HEADER_TITLE_LONG_LENGTH = 45


def _rgb(value: str, fallback: str = "#000000") -> RGBColor:
    color = value if isinstance(value, str) and value.startswith("#") and len(value) >= 7 else fallback
    return RGBColor.from_string(color[1:7])


def _add_shape(slide, shape_type, x, y, width, height, fill: str, line: str | None = None, radius=False):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
    else:
        shape.line.fill.background()
    return shape


def _add_text(
    slide, text: str, x: float, y: float, width: float, height: float,
    *, size: float = 18, color: str = "#111111", bold: bool = False,
    align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin: float = 0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _add_image(slide, image, x: float, y: float, width: float, height: float):
    stream = BytesIO(base64.b64decode(image.data))
    picture = slide.shapes.add_picture(stream, Inches(x), Inches(y), width=Inches(width), height=Inches(height))
    return picture


def _add_header(slide, page: dict, page_number: int, total: int, palette: dict) -> None:
    accent = palette["accent"]
    title = str(page.get("title") or "")
    if len(title) > HEADER_TITLE_LONG_LENGTH:
        title_size = HEADER_TITLE_LONG_SIZE
    elif len(title) > HEADER_TITLE_MEDIUM_LENGTH:
        title_size = HEADER_TITLE_MEDIUM_SIZE
    else:
        title_size = HEADER_TITLE_MAX_SIZE
    _add_text(slide, page.get("section_label", "PRESENTATION"), .68, .18, 2.4, .22,
              size=7.5, color=accent, bold=True)
    _add_text(slide, f"{page_number:02d} / {total:02d}", .65, .52, 1.05, .32,
              size=9, color=accent, bold=True, valign=MSO_ANCHOR.MIDDLE)
    _add_text(slide, title, 1.65, .36, 10.85, .62,
              size=title_size, color=palette.get("title_color", "#111111"), bold=True,
              valign=MSO_ANCHOR.MIDDLE)
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, .65, 1.10, 12.0, .035, accent)


def _add_intro(slide, text: str, y: float, palette: dict, x: float = .72, width: float = 11.9) -> float:
    if not text:
        return y
    _add_text(slide, text, x, y, width, .78, size=16, color=palette.get("body_color", "#374151"))
    return y + .82


def _add_numbered_cards(slide, bullets: list, y: float, palette: dict, *, x=.72, width=11.9, max_height=5.6) -> None:
    if not bullets:
        return
    gap = .10
    height = min(.76, max(.52, (max_height - gap * (len(bullets) - 1)) / len(bullets)))
    for index, bullet in enumerate(bullets[:6]):
        top = y + index * (height + gap)
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, width, height,
                   palette.get("card_bg", "#f9fafb"), palette.get("border_color", "#e5e7eb"))
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x + .18, top + (height - .32) / 2, .32, .32, palette["accent"])
        _add_text(slide, str(index + 1), x + .18, top + (height - .32) / 2, .32, .32,
                  size=9, color="#ffffff", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, bullet, x + .65, top + .08, width - .82, height - .12,
                  size=16, color=palette.get("body_color", "#374151"), valign=MSO_ANCHOR.MIDDLE)


def _render_cover(slide, page: dict, page_number: int, total: int, palette: dict, closing: bool = False) -> None:
    background = palette.get("slide_bg", "#ffffff")
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, background)
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, .12, 7.5, palette.get("cover_side_bar", palette["accent"]))
    # Keep decorative geometry inside the slide canvas. PowerPoint does not
    # clip off-canvas shapes consistently and QA tools correctly flag them as
    # overflow.
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 10.55, 0, 2.78, 2.78, palette.get("cover_deco", palette["accent2"]))
    color = palette.get("cover_text", palette.get("title_color", "#111111"))
    sub_color = palette.get("title_color", "#555555") if palette.get("cover_text_sub", "").startswith("rgba") else palette.get("cover_text_sub", "#555555")
    align = PP_ALIGN.CENTER if closing else PP_ALIGN.LEFT
    x, width = (1.65, 10.0) if closing else (1.05, 10.7)
    title_y = 2.35 if closing else 2.15
    if closing:
        _add_text(slide, "✦", 6.25, 1.45, .8, .6, size=24, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, page.get("title", ""), x, title_y, width, 1.55, size=44 if closing else 50,
              color=color, bold=True, align=align, valign=MSO_ANCHOR.MIDDLE)
    subtitle = page.get("subtitle") or ""
    content = page.get("content") or ""
    if subtitle:
        _add_text(slide, subtitle, x, title_y + 1.58, width, .82, size=22, color=sub_color, align=align)
    if content:
        _add_text(slide, content, x, title_y + 2.48, width, .82, size=15, color=sub_color, align=align)
    if not closing:
        for index, tag in enumerate((page.get("bullets") or [])[:5]):
            tag_x = 1.05 + index * 1.65
            _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, tag_x, 5.65, 1.48, .38, palette.get("tag_bg", palette["accent"]))
            _add_text(slide, tag, tag_x + .05, 5.68, 1.38, .30, size=9, color=palette.get("tag_text", "#ffffff"), bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, f"{page_number} / {total}", 11.7, 7.0, .9, .25, size=8, color=sub_color, align=PP_ALIGN.RIGHT)


def _render_stats(slide, page: dict, palette: dict) -> None:
    y = _add_intro(slide, page.get("content") or "", 1.32, palette)
    stats = (page.get("stats") or [])[:4]
    count = max(len(stats), 1)
    gap = .22
    card_width = (11.9 - gap * (count - 1)) / count
    for index, stat in enumerate(stats):
        x = .72 + index * (card_width + gap)
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y + .25, card_width, 3.75,
                   palette.get("card_bg", "#f9fafb"), palette.get("border_color", "#e5e7eb"))
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y + .25, card_width, .07, palette["accent"])
        _add_text(slide, stat.get("value", ""), x + .15, y + .75, card_width - .3, .8,
                  size=27, color=palette["accent"], bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, stat.get("label", ""), x + .15, y + 1.62, card_width - .3, .45,
                  size=16, color=palette.get("title_color", "#111111"), bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, stat.get("desc", ""), x + .22, y + 2.18, card_width - .44, 1.05,
                  size=14, color=palette.get("muted_color", "#6b7280"), align=PP_ALIGN.CENTER)


def _numeric_value(value: str) -> float:
    match = re.search(r"[-+]?\d[\d,.]*", str(value or ""))
    return float(match.group(0).replace(",", "")) if match else 0


def _render_data_chart(slide, page: dict, palette: dict) -> None:
    y = _add_intro(slide, page.get("content") or "", 1.35, palette)
    stats = (page.get("stats") or [])[:4]
    values = [abs(_numeric_value(stat.get("value", ""))) for stat in stats]
    max_value = max(values) if values else 1
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, .72, y + .12, 11.9, 4.25,
               palette.get("card_bg", "#f9fafb"), palette.get("border_color", "#e5e7eb"))
    for index, stat in enumerate(stats):
        top = y + .48 + index * .91
        _add_text(slide, stat.get("label", ""), 1.05, top, 2.15, .32, size=14,
                  color=palette.get("title_color", "#111111"), bold=True)
        track_width = 6.8
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 3.25, top + .02, track_width, .24,
                   palette.get("card_bg2", "#f1f3f5"))
        bar_width = max(.35, track_width * (values[index] / max_value)) if max_value else .35
        bar_color = palette["accent"] if index % 2 == 0 else palette["accent2"]
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 3.25, top + .02, bar_width, .24, bar_color)
        _add_text(slide, stat.get("value", ""), 10.25, top - .06, 1.55, .34, size=14,
                  color=bar_color, bold=True, align=PP_ALIGN.RIGHT)
        _add_text(slide, stat.get("desc", ""), 3.25, top + .34, 8.55, .32, size=12,
                  color=palette.get("muted_color", "#6b7280"))


def _render_quote(slide, page: dict, palette: dict) -> None:
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.15, 1.55, 11.0, 4.7,
               palette.get("card_bg2", "#f1f3f5"))
    _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 1.15, 1.55, .09, 4.7, palette["accent"])
    _add_text(slide, "“", 1.65, 1.72, .65, .65, size=40, color=palette["accent"], bold=True)
    _add_text(slide, page.get("quote") or "", 2.15, 2.25, 9.0, 1.65, size=23,
              color=palette.get("title_color", "#111111"), bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    _add_text(slide, page.get("content") or "", 2.2, 4.35, 8.9, .95, size=16,
              color=palette.get("body_color", "#374151"), align=PP_ALIGN.CENTER)


def _render_grid(slide, page: dict, palette: dict, spotlight: bool = False) -> None:
    y = 1.35
    if spotlight:
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, .72, y, 11.9, 1.25, palette.get("card_bg2", "#f1f3f5"))
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, .72, y, .08, 1.25, palette["accent"])
        _add_text(slide, page.get("content") or page.get("subtitle") or "", 1.05, y + .16, 11.1, .9,
                  size=19, color=palette.get("title_color", "#111111"), bold=True, valign=MSO_ANCHOR.MIDDLE)
        y += 1.55
    else:
        y = _add_intro(slide, page.get("content") or "", y, palette)
    bullets = (page.get("bullets") or [])[:6]
    card_width = 5.82
    for index, bullet in enumerate(bullets):
        column, row = index % 2, index // 2
        x = .72 + column * 6.08
        top = y + row * 1.18
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, card_width, .98,
                   palette.get("card_bg", "#f9fafb"), palette.get("border_color", "#e5e7eb"))
        _add_text(slide, f"{index + 1:02d}", x + .18, top + .14, .5, .25, size=9, color=palette["accent"], bold=True)
        _add_text(slide, bullet, x + .75, top + .10, 4.82, .76, size=15,
                  color=palette.get("body_color", "#374151"), valign=MSO_ANCHOR.MIDDLE)


def _render_timeline(slide, page: dict, palette: dict) -> None:
    y = _add_intro(slide, page.get("content") or "", 1.3, palette)
    bullets = (page.get("bullets") or [])[:6]
    for index, bullet in enumerate(bullets):
        top = y + index * .72
        if index < len(bullets) - 1:
            _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 1.02, top + .32, .025, .72, palette.get("border_color", "#e5e7eb"))
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, .91, top + .18, .25, .25, palette["accent"])
        _add_text(slide, f"{index + 1:02d}", 1.35, top + .10, .45, .35, size=10, color=palette["accent"], bold=True)
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.85, top, 10.3, .58,
                   palette.get("card_bg", "#f9fafb"), palette.get("border_color", "#e5e7eb"))
        _add_text(slide, bullet, 2.08, top + .05, 9.75, .46, size=15, color=palette.get("body_color", "#374151"), valign=MSO_ANCHOR.MIDDLE)


def _render_process(slide, page: dict, palette: dict) -> None:
    y = _add_intro(slide, page.get("content") or "", 1.35, palette)
    bullets = (page.get("bullets") or [])[:5]
    count = max(len(bullets), 1)
    gap = .18
    width = (11.9 - gap * (count - 1)) / count
    for index, bullet in enumerate(bullets):
        x = .72 + index * (width + gap)
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y + .25, width, 3.75,
                   palette.get("card_bg", "#f9fafb"), palette.get("border_color", "#e5e7eb"))
        _add_text(slide, f"STEP {index + 1:02d}", x + .2, y + .5, width - .4, .25,
                  size=8, color=palette["accent"], bold=True)
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x + .2, y + 1.05, .52, .52, palette["accent"])
        _add_text(slide, str(index + 1), x + .2, y + 1.05, .52, .52, size=12, color="#ffffff",
                  bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        _add_text(slide, bullet, x + .2, y + 1.85, width - .4, 1.45, size=14,
                  color=palette.get("body_color", "#374151"))


def _render_comparison(slide, page: dict, palette: dict) -> None:
    y = _add_intro(slide, page.get("content") or "", 1.35, palette)
    bullets = (page.get("bullets") or [])[:6]
    midpoint = (len(bullets) + 1) // 2
    for column, items in enumerate((bullets[:midpoint], bullets[midpoint:])):
        x = .72 + column * 6.08
        color = palette["accent"] if column == 0 else palette["accent2"]
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y + .15, 5.82, 4.15,
                   palette.get("card_bg", "#f9fafb"), palette.get("border_color", "#e5e7eb"))
        _add_text(slide, "A" if column == 0 else "B", x + .3, y + .42, 5.2, .4,
                  size=17, color=color, bold=True)
        _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + .3, y + .93, 5.2, .025, color)
        for index, bullet in enumerate(items):
            top = y + 1.25 + index * .82
            _add_text(slide, "◆", x + .32, top, .25, .25, size=8, color=color)
            _add_text(slide, bullet, x + .7, top - .03, 4.7, .62, size=14,
                      color=palette.get("body_color", "#374151"), valign=MSO_ANCHOR.MIDDLE)


def _render_content(slide, page: dict, palette: dict, image_map: dict) -> None:
    image = image_map.get(page.get("image_index"))
    image_position = page.get("image_position") or "right"
    if image and image_position in {"left", "right"}:
        image_x = .72 if image_position == "left" else 8.05
        text_x = 4.95 if image_position == "left" else .72
        _add_image(slide, image, image_x, 1.45, 4.55, 4.85)
        caption = page.get("image_caption") or ""
        if caption:
            _add_text(slide, caption, image_x, 6.38, 4.55, .28, size=8, color=palette.get("muted_color", "#6b7280"), align=PP_ALIGN.CENTER)
        y = _add_intro(slide, page.get("content") or "", 1.45, palette, text_x, 4.55)
        _add_numbered_cards(slide, page.get("bullets") or [], y, palette, x=text_x, width=4.55, max_height=5.0 - (y - 1.45))
        return
    y = 1.35
    if image and image_position == "full":
        _add_image(slide, image, .72, y, 11.9, 2.3)
        y += 2.48
    y = _add_intro(slide, page.get("content") or page.get("subtitle") or "", y, palette)
    _add_numbered_cards(slide, page.get("bullets") or [], y, palette, max_height=6.55 - y)


def build_pptx(page_data: dict, images: list, palette: dict, output_path: Path) -> None:
    """Build a native, editable 16:9 presentation from structured slide data."""
    page_data = prepare_presentation(page_data)
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    blank_layout = presentation.slide_layouts[6]
    pages = page_data.get("pages", [])
    image_map = {image.index: image for image in images}
    background = palette.get("slide_bg", "#ffffff")

    for index, page in enumerate(pages):
        slide = presentation.slides.add_slide(blank_layout)
        layout = page.get("layout", "content")
        if index == 0:
            layout = "cover"
        elif index == len(pages) - 1:
            layout = "closing"
        if layout in {"cover", "closing"}:
            _render_cover(slide, page, index + 1, len(pages), palette, layout == "closing")
        else:
            _add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, background)
            _add_header(slide, page, index + 1, len(pages), palette)
            if layout == "stats" and page.get("stats"):
                _render_stats(slide, page, palette)
            elif layout == "data_chart" and page.get("stats"):
                _render_data_chart(slide, page, palette)
            elif layout == "quote":
                _render_quote(slide, page, palette)
            elif layout == "timeline" and page.get("bullets"):
                _render_timeline(slide, page, palette)
            elif layout == "process" and page.get("bullets"):
                _render_process(slide, page, palette)
            elif layout == "comparison" and page.get("bullets"):
                _render_comparison(slide, page, palette)
            elif layout in {"card_grid", "spotlight"}:
                _render_grid(slide, page, palette, layout == "spotlight")
            else:
                _render_content(slide, page, palette, image_map)
        notes = page.get("speaker_notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    presentation.save(str(output_path))
