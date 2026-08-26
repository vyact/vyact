"""
document_parser.py – 파일 파싱 + 청크 분할
지원 형식: pdf, docx, xlsx, pptx, txt, html/htm, md
chunk_type: paragraph | table | code | heading | caption
"""
import re
from collections import Counter
from pathlib import Path
from typing import Iterable, NamedTuple

import pypdfium2 as pdfium

from logger import get_logger
from services.runtime_settings import get_runtime_settings

logger = get_logger(__name__)

CHUNK_SIZE = 1200
CHUNK_OVERLAP = 150
MIN_STANDALONE_CHUNK_LENGTH = 30
PDF_VECTOR_PATH_FAST_FALLBACK_THRESHOLD = 20_000


def _chunk_settings() -> tuple[int, int]:
    settings = get_runtime_settings()
    return settings["document_chunk_size"], settings["document_chunk_overlap"]


class Chunk(NamedTuple):
    text: str
    chunk_type: str          # paragraph | table | code | heading | caption
    heading_path: list[str] = []   # 소속 heading 경로 ["1. 서론", "1.1 배경"]
    page_number: int | None = None  # PDF 페이지 번호 (1-based)


def _split_lines_with_context(lines: Iterable[str], *, context_lines: list[str], chunk_type: str,
                              heading_path: list[str] | None = None, page_number: int | None = None) -> list[Chunk]:
    """행 경계를 지키며 모든 구조화 청크에 최대 길이를 적용한다."""
    chunk_size, _ = _chunk_settings()
    context = [line for line in context_lines if line.strip()]
    # 문맥 자체가 제한을 넘으면 행을 보존한 채 반복할 수 없다. 이 경우에는
    # 문맥을 별도 청크로 내보내고, 뒤따르는 행은 제한 내에서 분할한다.
    if context and len("\n".join(context)) >= chunk_size - 1:
        context_chunks = [
            Chunk(text=part, chunk_type=chunk_type, heading_path=list(heading_path or []), page_number=page_number)
            for part in split_chunks("\n".join(context), chunk_size=chunk_size, overlap=0)
        ]
        return [
            *context_chunks,
            *_split_lines_with_context(lines, context_lines=[], chunk_type=chunk_type,
                                       heading_path=heading_path, page_number=page_number),
        ]
    prefix_length = len("\n".join(context)) + (1 if context else 0)
    output: list[Chunk] = []
    current: list[str] = []
    current_length = prefix_length

    def emit() -> None:
        if current:
            output.append(Chunk(text="\n".join([*context, *current]).strip(), chunk_type=chunk_type,
                                heading_path=list(heading_path or []), page_number=page_number))
            current.clear()

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if current and current_length + len(line) + 1 > chunk_size:
            emit()
            current_length = prefix_length
        if not current and prefix_length + len(line) > chunk_size:
            available = max(1, chunk_size - prefix_length)
            for part in split_chunks(line, chunk_size=available, overlap=0):
                output.append(Chunk(text="\n".join([*context, part]).strip(), chunk_type=chunk_type,
                                    heading_path=list(heading_path or []), page_number=page_number))
            continue
        current.append(line)
        current_length += len(line) + 1
    emit()
    return output


def _table_chunks(rows: list[str], *, context_lines: list[str] | None = None,
                  heading_path: list[str] | None = None, page_number: int | None = None) -> list[Chunk]:
    """표는 행 경계에서 분할하고 제목·헤더를 모든 분할 청크에 반복한다."""
    if not rows:
        return []

    def is_section_row(row: str) -> bool:
        """병합 셀로 만든 짧은 섹션 제목 행인지 판별한다.

        문서 작성 도구는 서로 다른 논리 표를 하나의 외곽 표에 배치하는 경우가
        많다. 한 칸만 값이 있는 짧은 행은 새 표의 제목으로 취급해야 이전 표의
        헤더가 다음 섹션 데이터에 붙지 않는다.
        """
        cells = [cell.strip() for cell in row.split("|") if cell.strip()]
        return len(cells) == 1 and len(cells[0]) <= 80 and not re.match(r"^[●■▪•-]", cells[0])

    def split_sections(table_rows: list[str]) -> list[list[str]]:
        sections: list[list[str]] = []
        current: list[str] = []
        for row in table_rows:
            if current and is_section_row(row):
                sections.append(current)
                current = [row]
            else:
                current.append(row)
        if current:
            sections.append(current)
        return sections

    sections = split_sections(rows)
    if len(sections) > 1:
        result: list[Chunk] = []
        for section_index, section in enumerate(sections):
            section_context = list(context_lines or [])
            if section_index:
                section_context.append(section[0])
                section = section[1:]
            if section:
                result.extend(_table_chunks(section, context_lines=section_context,
                                            heading_path=heading_path, page_number=page_number))
            elif section_context:
                result.append(Chunk(text="\n".join(section_context), chunk_type="table",
                                    heading_path=list(heading_path or []), page_number=page_number))
        return result

    # 한 행짜리 표는 첫 행이 곧 전체 내용이다. 이를 헤더와 데이터 양쪽에 넣으면
    # 긴 셀을 분할할 때 원문 전체가 매번 중복된다.
    if len(rows) == 1:
        return _split_lines_with_context(rows, context_lines=list(context_lines or []),
                                         chunk_type="table", heading_path=heading_path, page_number=page_number)
    return _split_lines_with_context(rows[1:], context_lines=[*(context_lines or []), rows[0]],
                                     chunk_type="table", heading_path=heading_path, page_number=page_number)


# ─────────────────────────────
# 파서들 (텍스트만 반환 — 내부용)
# ─────────────────────────────

def _words_to_pdf_text(words: list[dict]) -> str:
    """좌표가 있는 PDF 단어를 위에서 아래, 왼쪽에서 오른쪽 순으로 재조립한다."""
    if not words:
        return ""
    lines: list[list[dict]] = []
    for word in sorted(words, key=lambda item: (item["top"], item["x0"])):
        if not lines or abs(lines[-1][0]["top"] - word["top"]) > 3:
            lines.append([word])
        else:
            lines[-1].append(word)
    output: list[str] = []
    for index, line in enumerate(lines):
        output.append(" ".join(word["text"] for word in sorted(line, key=lambda item: item["x0"])))
        if index + 1 < len(lines):
            line_bottom = max(word["bottom"] for word in line)
            next_top = lines[index + 1][0]["top"]
            if next_top - line_bottom > 6:
                output.append("")
    return "\n".join(output)


def _extract_pdf_body_text(
    page,
    table_bboxes: list[tuple],
    repeated_margin_words: set[str] | None = None,
    page_words: list[dict] | None = None,
) -> str:
    """표를 제외한 PDF 본문을 읽기 순서대로 추출한다.

    중앙 여백이 뚜렷한 페이지는 좌측 열 전체 후 우측 열 전체를 읽는다. 단일 열
    문서나 넓은 표제/초록이 있는 페이지는 기존의 위→아래 추출을 유지한다.
    """
    try:
        words = page_words if page_words is not None else page.extract_words()
        if table_bboxes:
            words = [
                word for word in words
                if not any(
                    word["x0"] < bbox[2]
                    and word["x1"] > bbox[0]
                    and word["top"] < bbox[3]
                    and word["bottom"] > bbox[1]
                    for bbox in table_bboxes
                )
            ]
        if repeated_margin_words:
            margin_height = min(40, page.height * 0.08)
            words = [
                word for word in words
                if not (
                    word["top"] <= margin_height or word["bottom"] >= page.height - margin_height
                ) or (
                    re.sub(r"\W+", "", word["text"]).lower() not in repeated_margin_words
                    and not (
                        word["bottom"] >= page.height - margin_height
                        and re.fullmatch(r"\d{1,3}", word["text"].strip())
                    )
                )
            ]
        if not words:
            return ""

        page_midpoint = page.width / 2
        gutter = page.width * 0.01
        # x0/x1 기준은 한 줄 전체가 하나의 word로 추출된 경우 양쪽 열에 같은
        # 텍스트를 넣는다. 단어 중심점으로 양쪽 열을 상호 배타적으로 나눈다.
        left_words = [
            word for word in words
            if (word["x0"] + word["x1"]) / 2 < page_midpoint - gutter
        ]
        right_words = [
            word for word in words
            if (word["x0"] + word["x1"]) / 2 > page_midpoint + gutter
        ]
        center_words = [
            word for word in words
            if word["x0"] < page_midpoint - gutter and word["x1"] > page_midpoint + gutter
        ]
        has_two_columns = (
            len(left_words) >= 30 and len(right_words) >= 30
            and len(center_words) * 4 < min(len(left_words), len(right_words))
        )
        if has_two_columns:
            left_text = _words_to_pdf_text(left_words)
            right_text = _words_to_pdf_text(right_words)
            return "\n".join(part for part in (left_text, right_text) if part.strip())
        return _words_to_pdf_text(words)
    except Exception:
        return page.extract_text() or ""

def _parse_pdf_content(path: Path) -> tuple[list[Chunk], str]:
    """한 번 연 PDF에서 typed 청크와 원문 텍스트를 함께 추출한다."""
    try:
        import pdfplumber
    except ImportError:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        text = "\n\n".join(p.extract_text() or "" for p in reader.pages)
        return _text_to_chunks(text), text

    chunks: list[Chunk] = []
    original_pages: list[str] = []
    current_headings: list[str] = []  # 현재 heading 경로 스택

    pdfium_document = pdfium.PdfDocument(path)
    complex_page_text: dict[int, str] = {}
    try:
        for page_index in range(len(pdfium_document)):
            pdfium_page = pdfium_document[page_index]
            path_count = sum(1 for obj in pdfium_page.get_objects() if obj.type == 2)
            if path_count >= PDF_VECTOR_PATH_FAST_FALLBACK_THRESHOLD:
                text_page = pdfium_page.get_textpage()
                try:
                    complex_page_text[page_index] = text_page.get_text_range()
                finally:
                    text_page.close()
            pdfium_page.close()
    finally:
        pdfium_document.close()

    with pdfplumber.open(path) as pdf:
        margin_counts: Counter[str] = Counter()
        page_words_list: list[list[dict] | None] = []
        for page_index, page in enumerate(pdf.pages):
            if page_index in complex_page_text:
                page_words_list.append(None)
                continue
            page_words = page.extract_words()
            page_words_list.append(page_words)
            margin_height = min(40, page.height * 0.08)
            page_margin_words = {
                re.sub(r"\W+", "", word["text"]).lower()
                for word in page_words
                if word["top"] <= margin_height or word["bottom"] >= page.height - margin_height
            }
            margin_counts.update(word for word in page_margin_words if len(word) >= 3)
        repeated_margin_words = {word for word, count in margin_counts.items() if count >= 2}

        for page_num, (page, page_words) in enumerate(zip(pdf.pages, page_words_list), start=1):
            if page_words is None:
                page_text = complex_page_text[page_num - 1]
                if page_text.strip():
                    original_pages.append(page_text.strip())
                    chunks.extend(
                        _classify_text_lines_with_context(page_text, current_headings, page_num)
                    )
                continue

            # 복잡한 벡터 페이지에서 extract_text로 다시 해석하지 않고, 검색 본문과
            # 동일한 좌표 단어를 원문에도 재사용한다.
            original_page_text = _words_to_pdf_text(page_words)
            if original_page_text.strip():
                original_pages.append(original_page_text.strip())

            # 표 탐지는 한 번만 수행하고 같은 결과에서 셀과 bbox를 함께 사용한다.
            found_tables = page.find_tables() if hasattr(page, "find_tables") else []
            tables = [table.extract() for table in found_tables]
            table_bboxes = [table.bbox for table in found_tables]

            for table in tables:
                rows = []
                for row in table:
                    cells = [str(c).strip() if c else "" for c in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
                chunks.extend(_table_chunks(rows, heading_path=current_headings, page_number=page_num))

            # 2) 표 영역을 제외하고 단일/다단 레이아웃에 맞춰 본문을 추출
            page_text = _extract_pdf_body_text(
                page,
                table_bboxes,
                repeated_margin_words,
                page_words=page_words,
            )

            if not page_text:
                continue

            # 3) 줄 단위 분석 → heading 경로 추적
            for raw_chunk in _classify_text_lines_with_context(page_text, current_headings, page_num):
                chunks.append(raw_chunk)

    return chunks, "\n\n".join(original_pages)


def _parse_pdf_chunks(path: Path) -> list[Chunk]:
    """pdfplumber로 표/텍스트/코드/제목 구분 청크 생성"""
    return _parse_pdf_content(path)[0]


def _classify_text_lines_with_context(
        text: str,
        heading_stack: list[str],
        page_num: int | None = None,
) -> list[Chunk]:
    """텍스트를 줄 단위로 분석 + heading 경로 추적"""
    chunks: list[Chunk] = []
    lines = text.split("\n")

    code_buf: list[str] = []
    caption_buf: list[str] = []
    para_buf: list[str] = []

    CODE_PATTERNS = re.compile(
        r"^(def |class |public |private |protected |function |const |let |var |{|}|#include|package |@)"
    )
    CAPTION_PATTERN = re.compile(
        r"^(?:fig(?:ure)?\.?|table)\s*\d+(?:\||[.:]\s+|\s+)",
        re.IGNORECASE,
    )
    HEADING_PATTERN = re.compile(
        r"^("
        r"#{1,4}\s"                       # 마크다운 헤딩: ## 제목
        r"|제\d+[장절]"                    # 한국어 장/절: 제1장
        r")"
    )

    def flush_para():
        t = " ".join(para_buf).strip()
        if t:
            chunks.extend(
                Chunk(
                    text=part, chunk_type="paragraph",
                    heading_path=list(heading_stack),
                    page_number=page_num,
                )
                for part in split_chunks(t)
            )
        para_buf.clear()

    def flush_code():
        t = "\n".join(code_buf).strip()
        if t:
            chunks.append(Chunk(
                text=t, chunk_type="code",
                heading_path=list(heading_stack),
                page_number=page_num,
            ))
        code_buf.clear()

    def flush_caption():
        t = " ".join(caption_buf).strip()
        if t:
            chunks.append(Chunk(
                text=t, chunk_type="caption",
                heading_path=list(heading_stack),
                page_number=page_num,
            ))
        caption_buf.clear()

    pending_heading: str | None = None

    for line in lines:
        stripped = line.strip()
        if not stripped:
            if code_buf:
                flush_code()
            elif caption_buf:
                flush_caption()
            # PDF line spacing frequently appears as an empty line. Keep prose
            # together until a structural boundary or the end of the page.
            continue

        if caption_buf:
            caption_buf.append(stripped)
        elif CAPTION_PATTERN.match(stripped):
            if code_buf:
                flush_code()
            if para_buf:
                flush_para()
            caption_buf.append(stripped)
        elif CODE_PATTERNS.match(stripped):
            if para_buf:
                flush_para()
            code_buf.append(stripped)
        elif HEADING_PATTERN.match(stripped) and len(stripped) < 80:
            if code_buf:
                flush_code()
            if para_buf:
                flush_para()
            # heading_stack 업데이트: 간단히 최대 3레벨 유지
            heading_stack.clear()
            heading_stack.append(stripped)
            chunks.append(Chunk(
                text=stripped, chunk_type="heading",
                heading_path=list(heading_stack),
                page_number=page_num,
            ))
        else:
            if code_buf:
                flush_code()
            if pending_heading:
                para_buf.append(pending_heading)
                pending_heading = None
            para_buf.append(stripped)

    if pending_heading:
        chunks.append(Chunk(
            text=pending_heading, chunk_type="heading",
            heading_path=list(heading_stack),
            page_number=page_num,
        ))
    flush_code()
    flush_caption()
    flush_para()
    return chunks


def _text_to_chunks(text: str) -> list[Chunk]:
    """일반 텍스트 → paragraph 청크 리스트"""
    return [Chunk(text=c, chunk_type="paragraph") for c in split_chunks(text) if c.strip()]


def _parse_docx_chunks(path: Path) -> list[Chunk]:
    """DOCX의 문단·표 순서와 병합 셀을 보존해 청킹한다."""
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    doc = Document(str(path))
    chunks: list[Chunk] = []
    para_buf: list[str] = []
    heading_stack: list[str] = []

    HEADING_STYLES = {"heading 1", "heading 2", "heading 3", "heading 4", "제목 1", "제목 2"}

    def flush_para():
        t = " ".join(para_buf).strip()
        if t:
            for c in split_chunks(t):
                chunks.append(Chunk(text=c, chunk_type="paragraph", heading_path=list(heading_stack)))
        para_buf.clear()

    def append_table(table: Table) -> None:
        rows = []
        for row in table.rows:
            seen_cells = set()
            values = []
            for cell in row.cells:
                cell_id = id(cell._tc)
                if cell_id in seen_cells:
                    continue
                seen_cells.add(cell_id)
                if cell.text.strip():
                    values.append(cell.text.strip())
            if values:
                rows.append(" | ".join(values))
        chunks.extend(_table_chunks(rows, heading_path=heading_stack))

    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_Tbl):
            flush_para()
            append_table(Table(child, doc))
            continue
        if not isinstance(child, CT_P):
            continue
        para = Paragraph(child, doc)
        t = para.text.strip()
        if not t:
            continue
        style_name = para.style.name if para.style else ""
        style = style_name.lower()
        is_heading = style in HEADING_STYLES or style.startswith("heading")
        if is_heading:
            flush_para()
            # 레벨 파악 (Heading 1 → 레벨 1)
            level = 1
            try:
                level = int(style_name.split()[-1])
            except (ValueError, IndexError):
                pass
            # 상위 레벨 이후 항목 제거하고 현재 추가
            heading_stack[level - 1:] = [t]
            chunks.append(Chunk(text=t, chunk_type="heading", heading_path=list(heading_stack)))
        else:
            para_buf.append(t)

    flush_para()

    return chunks


def _parse_pdf(path: Path) -> str:
    return _parse_pdf_content(path)[1]


def _parse_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _parse_xlsx(path: Path) -> str:
    """XLSX typed chunks as complete plain text."""
    chunks = _parse_xlsx_chunks(path)
    return "\n\n".join(c.text for c in chunks)


def _parse_xlsx_chunks(path: Path) -> list["Chunk"]:
    """XLSX를 시트·행·헤더 문맥을 유지한 표 청크로 분할한다."""

    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=False, data_only=False)
    chunks: list[Chunk] = []

    for sheet in wb.worksheets:
        rows: list[str] = []
        for row_index, row in enumerate(sheet.iter_rows(), start=1):
            cells = [str(cell.value).strip() for cell in row if cell.value is not None and str(cell.value).strip()]
            if cells:
                rows.append(f"[행: {row_index}] " + " | ".join(cells))

        if not rows:
            continue

        header = rows[0]
        chunks.extend(_split_lines_with_context(
            rows[1:] or [header],
            context_lines=[f"[시트: {sheet.title}]", header],
            chunk_type="table",
        ))

    return chunks


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append(f"[슬라이드 {i}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _parse_pptx_chunks(path: Path) -> list[Chunk]:
    """슬라이드 경계를 넘지 않고 제목을 분할 청크의 문맥으로 유지한다."""
    from pptx import Presentation

    presentation = Presentation(str(path))
    chunks: list[Chunk] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        shapes = sorted((shape for shape in slide.shapes if shape.has_text_frame), key=lambda shape: (shape.top, shape.left))
        lines = [paragraph.text.strip() for shape in shapes for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
        if not lines:
            continue
        title = lines[0]
        chunks.extend(_split_lines_with_context(lines[1:] or [title], context_lines=[f"[슬라이드 {slide_number}]", title],
                                                chunk_type="paragraph", heading_path=[title], page_number=slide_number))
    return chunks


def _parse_txt(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp949", "euc-kr"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _parse_html(path: Path) -> str:
    from bs4 import BeautifulSoup
    raw = _parse_txt(path)
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _parse_html_chunks(path: Path) -> list["Chunk"]:
    """HTML → heading / table / code / paragraph 청크 분할"""
    from bs4 import BeautifulSoup, Tag
    raw = _parse_txt(path)
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    chunks: list[Chunk] = []
    para_buf: list[str] = []
    heading_stack: list[str] = []

    def flush_para():
        t = " ".join(para_buf).strip()
        if t:
            for c in split_chunks(t):
                chunks.append(Chunk(text=c, chunk_type="paragraph", heading_path=list(heading_stack)))
        para_buf.clear()

    root = soup.body if soup.body else soup
    for el in root.children:
        if not isinstance(el, Tag):
            continue

        if el.name in ("h1", "h2", "h3", "h4", "h5", "h6"):
            flush_para()
            t = el.get_text(strip=True)
            if t:
                level = int(el.name[1])
                heading_stack[level - 1:] = [t]
                chunks.append(Chunk(text=t, chunk_type="heading", heading_path=list(heading_stack)))

        elif el.name == "table":
            flush_para()
            rows = []
            for tr in el.find_all("tr"):
                cells = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
                cells = [c for c in cells if c]
                if cells:
                    rows.append(" | ".join(cells))
            chunks.extend(_table_chunks(rows, heading_path=heading_stack))

        elif el.name in ("pre", "code"):
            flush_para()
            t = el.get_text(strip=True)
            if t:
                chunks.append(Chunk(text=t, chunk_type="code", heading_path=list(heading_stack)))

        else:
            t = el.get_text(separator=" ", strip=True)
            if t:
                para_buf.append(t)

    flush_para()
    return chunks


def _parse_md(path: Path) -> str:
    return _parse_txt(path)


def _parse_md_chunks(path: Path) -> list["Chunk"]:
    """마크다운 → heading / table / code / paragraph 청크 분할"""
    text = _parse_txt(path)
    chunks: list[Chunk] = []

    code_buf: list[str] = []
    table_buf: list[str] = []
    html_buf: list[str] = []
    para_buf: list[str] = []
    in_code_block = False
    heading_stack: list[str] = []  # 인덱스 = 레벨-1

    def flush_para():
        t = " ".join(para_buf).strip()
        if t:
            for c in split_chunks(t):
                chunks.append(Chunk(text=c, chunk_type="paragraph", heading_path=list(heading_stack)))
        para_buf.clear()

    def flush_table():
        rows = [r for r in table_buf if not re.match(r"^\|[\s\-|:]+\|$", r.strip())]
        if rows:
            chunks.extend(_table_chunks(rows, heading_path=heading_stack))
        table_buf.clear()

    def flush_code():
        t = "\n".join(code_buf).strip()
        if t:
            chunks.append(Chunk(text=t, chunk_type="code", heading_path=list(heading_stack)))
        code_buf.clear()

    def flush_html():
        if not html_buf:
            return
        from bs4 import BeautifulSoup

        soup = BeautifulSoup("\n".join(html_buf), "html.parser")
        for table in soup.find_all("table"):
            rows = []
            for tr in table.find_all("tr"):
                cells = [cell.get_text(" ", strip=True) for cell in tr.find_all(["th", "td"])]
                if any(cells):
                    rows.append(" | ".join(cells))
            chunks.extend(_table_chunks(rows, heading_path=heading_stack))
            table.decompose()
        text = soup.get_text(" ", strip=True)
        if text:
            for chunk in split_chunks(text):
                chunks.append(Chunk(text=chunk, chunk_type="paragraph", heading_path=list(heading_stack)))
        html_buf.clear()

    for line in text.split("\n"):
        # Markdown 수평선은 문서 구조를 나누는 표식일 뿐 검색 대상이 아니다.
        if re.fullmatch(r"\s{0,3}([-*_])(?:\s*\1){2,}\s*", line):
            flush_para()
            flush_table()
            continue
        if line.strip().startswith("```"):
            if in_code_block:
                flush_code()
                in_code_block = False
            else:
                flush_para()
                flush_table()
                in_code_block = True
            continue

        if in_code_block:
            code_buf.append(line)
            continue

        is_html_block = bool(re.match(
            r"\s*</?(?:section|article|div|table|thead|tbody|tfoot|tr|td|th|p|h[1-6]|ul|ol|li)\b",
            line,
            re.IGNORECASE,
        ))
        if html_buf:
            html_buf.append(line)
            if not line.strip():
                flush_html()
            continue
        if is_html_block:
            flush_para()
            flush_table()
            html_buf.append(line)
            continue

        if line.strip().startswith("|"):
            if para_buf:
                flush_para()
            table_buf.append(line.strip())
            continue
        else:
            if table_buf:
                flush_table()

        heading_match = re.match(r"^(#{1,6})\s+(.+)", line)
        if heading_match:
            flush_para()
            level = len(heading_match.group(1))
            h_text = line.strip()
            # 해당 레벨 이후 제거 후 현재 추가
            heading_stack[level - 1:] = [h_text]
            chunks.append(Chunk(text=h_text, chunk_type="heading", heading_path=list(heading_stack)))
            continue

        if not line.strip():
            if para_buf:
                flush_para()
            continue

        para_buf.append(line.strip())

    if in_code_block:
        flush_code()
    flush_html()
    flush_table()
    flush_para()

    return chunks


# ─────────────────────────────
# 청크 분할
# ─────────────────────────────

def split_chunks(text: str, chunk_size: int | None = None, overlap: int | None = None) -> list[str]:
    """문장 경계 기준 청크 분할 — 한국어 종결어미 포함"""
    default_size, default_overlap = _chunk_settings()
    chunk_size = default_size if chunk_size is None else chunk_size
    overlap = default_overlap if overlap is None else overlap
    text = re.sub(r"\n{3,}", "\n\n", text.strip())
    if len(text) <= chunk_size:
        return [text] if text else []

    # 한국어 종결어미 + 영문 문장부호 + 단락 경계
    sentences = re.split(
        r"(?<=[.!?。])\s+"           # 영문 문장부호 뒤 공백
        r"|(?<=다\.)\s+"             # ~다.
        r"|(?<=요\.)\s+"             # ~요.
        r"|(?<=죠\.)\s+"             # ~죠.
        r"|(?<=까\.)\s+"             # ~까.
        r"|(?<=니다\.)\s+"           # ~니다.
        r"|(?<=군요\.)\s+"           # ~군요.
        r"|(?<=네요\.)\s+"           # ~네요.
        r"|\n{2,}",                  # 단락 경계
        text
    )
    sentences = [s.strip() for s in sentences if s.strip()]

    chunks = []
    current = []
    current_len = 0

    for sent in sentences:
        sent_len = len(sent)
        if sent_len > chunk_size:
            if current:
                chunks.append(" ".join(current))
                current, current_len = [], 0
            # 문장 내부에 자연스러운 경계가 없을 때만 최후 수단으로 길이를 강제한다.
            for start in range(0, sent_len, chunk_size):
                chunks.append(sent[start:start + chunk_size].strip())
            continue
        if current and current_len + 1 + sent_len > chunk_size:
            chunk_text = " ".join(current)
            chunks.append(chunk_text)
            # overlap: 문자 슬라이싱 대신 마지막 문장 단위로
            overlap_sents = []
            overlap_len = 0
            for s in reversed(current):
                if overlap_len + len(s) > overlap:
                    break
                overlap_sents.insert(0, s)
                overlap_len += len(s)
            # 겹침 문장과 다음 문장이 함께 상한을 넘지 않게 한다.
            while overlap_sents and overlap_len + 1 + sent_len > chunk_size:
                overlap_len -= len(overlap_sents.pop(0))
            current = overlap_sents + [sent]
            current_len = len(" ".join(current))
        else:
            current.append(sent)
            current_len += sent_len + (1 if len(current) > 1 else 0)

    if current:
        chunks.append(" ".join(current))

    return [c for c in chunks if c.strip()]


# ─────────────────────────────
# 메인 파서
# ─────────────────────────────

PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".pptx": _parse_pptx,
    ".txt": _parse_txt,
    ".html": _parse_html,
    ".htm": _parse_html,
    ".md": _parse_md,
}

# chunk_type 지원 파서
CHUNK_PARSERS = {
    ".pdf":  _parse_pdf_chunks,
    ".docx": _parse_docx_chunks,
    ".xlsx": _parse_xlsx_chunks,
    ".html": _parse_html_chunks,
    ".htm":  _parse_html_chunks,
    ".md":   _parse_md_chunks,
    ".pptx": _parse_pptx_chunks,
}


def parse_file(path: Path) -> str:
    """파일 파싱 → 전체 텍스트 반환"""
    ext = path.suffix.lower()
    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(f"지원하지 않는 형식: {ext}")
    try:
        text = parser(path)
        logger.info("파싱 완료: %s (%d자)", path.name, len(text))
        return text
    except Exception as e:
        logger.error("파싱 실패 [%s]: %s", path.name, e)
        raise


def _normalize_typed_chunks(path: Path, chunks: list[Chunk]) -> list[Chunk]:
    """모든 typed 블록에 최대 길이를 적용하면서 구조 메타데이터를 보존한다."""
    result: list[Chunk] = []
    chunk_size, _ = _chunk_settings()
    for chunk in chunks:
        if len(chunk.text) > chunk_size:
            if chunk.chunk_type == "table":
                rows = [line for line in chunk.text.splitlines() if line.strip()]
                context_lines = []
                if rows and rows[0].startswith("[시트:"):
                    context_lines.append(rows.pop(0))
                result.extend(_table_chunks(rows, context_lines=context_lines, heading_path=chunk.heading_path, page_number=chunk.page_number))
            elif chunk.chunk_type == "code":
                result.extend(_split_lines_with_context(chunk.text.splitlines(), context_lines=[], chunk_type="code", heading_path=chunk.heading_path, page_number=chunk.page_number))
            else:
                for sub in split_chunks(chunk.text):
                    result.append(Chunk(text=sub, chunk_type=chunk.chunk_type, heading_path=chunk.heading_path, page_number=chunk.page_number))
        else:
            result.append(chunk)
    consolidated = _consolidate_short_chunks(result, chunk_size)
    logger.info("typed 청크 분할: %s → %d개", path.name, len(consolidated))
    return consolidated


def _consolidate_short_chunks(chunks: list[Chunk], chunk_size: int) -> list[Chunk]:
    """검색 가치가 낮은 짧은 청크를 제거하거나 같은 페이지의 이웃에 병합한다."""
    pending = list(chunks)
    output: list[Chunk] = []

    def can_merge(left: Chunk, right: Chunk) -> bool:
        return left.page_number == right.page_number and len(left.text) + len(right.text) + 1 <= chunk_size

    def merged(target: Chunk, addition: Chunk, prepend: bool) -> Chunk:
        parts = (addition.text, target.text) if prepend else (target.text, addition.text)
        return Chunk(
            text="\n".join(part.strip() for part in parts if part.strip()),
            chunk_type=target.chunk_type,
            heading_path=target.heading_path,
            page_number=target.page_number,
        )

    for index, chunk in enumerate(pending):
        text = chunk.text.strip()
        if not text:
            continue
        chunk = chunk._replace(text=text)
        if len(text) >= MIN_STANDALONE_CHUNK_LENGTH:
            output.append(chunk)
            continue

        # 페이지 번호처럼 문자 없이 숫자·기호만 있는 독립 문단은 검색 가치가 없다.
        if chunk.chunk_type == "paragraph" and not any(character.isalpha() for character in text):
            continue

        previous = output[-1] if output else None
        next_chunk = pending[index + 1] if index + 1 < len(pending) else None
        previous_type_matches = previous is not None and (
            previous.chunk_type == chunk.chunk_type or chunk.chunk_type in {"heading", "caption"}
        )
        next_type_matches = next_chunk is not None and (
            next_chunk.chunk_type == chunk.chunk_type or chunk.chunk_type in {"heading", "caption"}
        )

        # 제목과 캡션은 뒤따르는 본문/표의 검색 문맥으로 사용하는 편이 자연스럽다.
        if next_type_matches and next_chunk is not None and can_merge(chunk, next_chunk):
            pending[index + 1] = merged(next_chunk, chunk, prepend=True)
        elif previous_type_matches and previous is not None and can_merge(previous, chunk):
            output[-1] = merged(previous, chunk, prepend=False)
        else:
            output.append(chunk)

    return output


def parse_file_to_typed_chunks(path: Path) -> list[Chunk]:
    """파일 파싱 → (text, chunk_type) 청크 리스트 반환"""
    ext = path.suffix.lower()
    if ext in CHUNK_PARSERS:
        try:
            return _normalize_typed_chunks(path, CHUNK_PARSERS[ext](path))
        except Exception as e:
            logger.warning("typed 파싱 실패, fallback: %s", e)

    # fallback: 일반 텍스트 파싱
    text = parse_file(path)
    return _text_to_chunks(text)


def parse_file_for_indexing(path: Path) -> tuple[list[Chunk], str]:
    """인덱싱용 typed 청크와 원문을 중복 파싱 없이 가능한 범위에서 함께 반환한다."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        chunks, original_text = _parse_pdf_content(path)
        return _normalize_typed_chunks(path, chunks), original_text

    typed_chunks = parse_file_to_typed_chunks(path)
    if ext == ".xlsx":
        return typed_chunks, "\n\n".join(chunk.text for chunk in typed_chunks)
    return typed_chunks, parse_file(path)
