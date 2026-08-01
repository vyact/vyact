"""
routers/pdf.py – PDF 생성 엔드포인트
HTML → Playwright → PDF 파이프라인
"""
import asyncio
import base64
import json
import re
import tempfile
import uuid
from datetime import datetime, timezone
from pathlib import Path

from fastapi import APIRouter
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches

from agent import query_llm, collect_llm_stream, save_conversation, get_model_name
from routers.deps import INSTALL_DIR, sse, load_config_async
from services.db import DOCUMENT_ORIGINALS_INDEX, DOC_CHUNKS_INDEX, get_es
from services.runtime_settings import (
    reset_request_temperature_override,
    set_request_temperature_override,
)
from logger import get_logger

logger = get_logger(__name__)

router = APIRouter()

FILES_DIR = INSTALL_DIR / "temp"
FILES_DIR.mkdir(parents=True, exist_ok=True)

from config import INSTALL_DIR as _INSTALL_DIR

IMAGES_DIR = _INSTALL_DIR / "uploads" / "images"
IMAGES_DIR.mkdir(parents=True, exist_ok=True)

PRESENTATION_CONTEXT_CHAR_LIMIT = 120_000
PRESENTATION_TEMPERATURE = 0.55

# ── 팔레트 ─────────────────────────────────────────────────────────────────────
STYLE_PALETTES = {
    # ── 화이트: 밝은 커버 + 오렌지 accent (소개서 스타일) ──
    "white": {
        "name": "클린 화이트",
        # 커버: 밝은 크림/오프화이트 배경
        "cover_bg": "linear-gradient(160deg, #ffffff 0%, #f8f4ef 50%, #f0ebe3 100%)",
        "cover_text": "#111111",          # 커버 제목 — 다크
        "cover_text_sub": "#555555",      # 커버 서브텍스트
        "cover_side_bar": "#f97316",      # 좌측 세로 바
        "cover_deco": "#f97316",          # 데코 원
        "accent":  "#f97316",             # 오렌지 — 번호 뱃지, 구분선, 태그
        "accent2": "#ea580c",             # 오렌지 다크 — 보조
        "title_color":  "#111111",
        "body_color":   "#374151",
        "muted_color":  "#6b7280",
        "card_bg":      "#f9fafb",
        "card_bg2":     "#f1f3f5",
        "border_color": "#e5e7eb",
        "tag_bg":       "#f97316",
        "tag_text":     "#ffffff",
    },
    # ── 다크: #212121 차콜 + 앰버 골드 accent ──
    "dark": {
        "name": "다크 모던",
        "cover_bg": "linear-gradient(160deg, #212121 0%, #252520 55%, #1c1c18 100%)",
        "cover_text": "#f1f5f9",
        "cover_text_sub": "rgba(241,245,249,0.75)",
        "cover_side_bar": "#f59e0b",
        "cover_deco": "#353530",
        "accent":  "#f59e0b",
        "accent2": "#fcd34d",
        "slide_bg":     "#212121",
        "title_color":  "#f1f5f9",
        "body_color":   "#cbd5e1",
        "muted_color":  "#94a3b8",
        "card_bg":      "#2a2a2a",
        "card_bg2":     "#333328",
        "border_color": "#3a3a32",
        "tag_bg":       "#f59e0b",
        "tag_text":     "#1a1a0a",
    },
}


class ImageMeta(BaseModel):
    index: int
    filename: str
    type: str
    data: str


class PdfGenerateRequest(BaseModel):
    prompt: str
    page_count: int = 8
    page_count_auto: bool = True
    language: str = "ko"
    style: str = "white"
    output_format: str = "pdf"
    articles: list = []
    images: list[ImageMeta] = []
    conv_id: str = ""
    messages: list = []
    reasoning: bool = False  # 추론 on/off (프론트 스위치 값)


# ── LLM 슬라이드 구조 생성 ──────────────────────────────────────────────────────
def _build_system_prompt(lang: str, page_count: int, page_count_auto: bool,
                         has_articles: bool, has_images: bool) -> str:
    language_names = {
        "ko": "Korean",
        "en": "English",
        "es": "Spanish",
        "fr": "French",
        "zh": "Simplified Chinese",
        "ja": "Japanese",
        "th": "Thai",
        "vi": "Vietnamese",
    }
    lang_str = language_names.get(lang, "English")

    if page_count_auto:
        slide_rule = "Determine the optimal number of pages based on content complexity. Minimum 6, maximum 15. ALWAYS make the first page 'cover' layout and the LAST page 'closing' layout — no exceptions."
    else:
        slide_rule = f"Generate exactly {page_count} pages. First MUST be 'cover' layout. Last MUST be 'closing' layout — no exceptions."

    article_note = """
- CRITICALLY IMPORTANT: Base content HEAVILY on the provided articles/documents.
- Extract specific facts, numbers, dates, names, and expert quotes. Never use generic filler.
- Every page must reference at least one concrete data point from the source material.""" if has_articles else ""

    image_note = """
- Analyze each image carefully. Assign image_index only when page content directly relates to the image.
- image_position: "right"/"left" for side-by-side layouts, "full" only for dramatic full-bleed pages.""" if has_images else ""

    # 언어별 문체
    if lang == "ko":
        tone_rule = """
- Korean tone: Use natural, complete sentences ending in "~합니다", "~입니다", "~됩니다".
  ✗ Avoid overly abbreviated noun-endings for body text.
  ✓ Correct: "매출은 전년 대비 23% 증가했습니다.", "핵심 성장 동력은 AI 반도체입니다."
- Bullets: concise noun-phrase or short sentence (under 20 words). Include specific numbers/names.
- Use professional Korean business writing. No informal endings."""
    elif lang == "en":
        tone_rule = """
- Use clear, professional English. Complete sentences for content. Concise noun phrases for bullets.
- Include specific numbers, names, and data points."""
    else:
        tone_rule = f"""
- Write every visible text field exclusively in natural, professional {lang_str}.
- Use complete sentences for content and concise, readable phrases for bullets.
- Include specific numbers, names, and data points without mixing in English except proper nouns or standard abbreviations."""

    return f"""You are an expert presentation designer and analyst. Create a highly professional, content-rich PDF document.

OUTPUT FORMAT — CRITICAL:
Respond ONLY with a single valid JSON object. No markdown, no explanation. Start with {{ end with }}.

JSON structure:
{{
  "presentation_title": "concise compelling title",
  "pages": [
    {{
      "index": 0,
      "layout": "cover|content|two_column|stats|quote|timeline|spotlight|card_grid|image_focus|closing",
      "title": "page title (concise, under 40 chars)",
      "subtitle": "supporting subtitle or null",
      "content": "2-4 sentence intro paragraph — always fill this for content layout",
      "bullets": ["specific insight with data or name"] or null,
      "stats": [{{"label": "metric name", "value": "precise number/figure", "desc": "1-2 sentence context"}}] or null,
      "quote": "verbatim or paraphrased notable quote" or null,
      "image_index": integer or null,
      "image_position": "right|left|full" or null,
      "image_caption": "descriptive caption" or null,
      "speaker_notes": "1-2 sentence presenter note"
    }}
  ]
}}

LAYOUT RULES (strictly follow):
- {slide_rule}
- cover   → title (main heading) + subtitle + content (2-3 sentences overview). bullets = tech/keyword tags (3-6 short words/phrases, no sentences).
- closing → LAST PAGE ONLY. layout MUST be "closing". title = conclusion heading. subtitle = one-line theme. content = 2-3 sentence concluding paragraph. bullets = null (NEVER add bullets). stats = null. This page uses a special full-page design — do NOT use content/stats layout for the last page.
- stats   → stats MUST have 3-4 items with real numbers. content = brief intro. bullets = null.
- quote   → quote MUST NOT be null. content = analysis/context paragraph. bullets = null.
- timeline → use 3-6 bullets as a chronological sequence, each beginning with a date, phase, or clear order marker. content = a short framing paragraph.
- spotlight → use one decisive insight or conclusion in content, with 2-4 supporting bullets. Assign an image when one strengthens the message.
- card_grid → use 4-6 independent bullets that can be read as distinct themes, actions, risks, or opportunities. content = a concise setup paragraph.
- image_focus → use only when an attached image can carry the page visually. Assign image_position "full" and write a short, impactful content paragraph plus up to 3 bullets.
- content → ALWAYS fill BOTH content (intro paragraph, 2-3 sentences) AND bullets (4-6 items, each under 20 words with specific data). NEVER leave both null.
- two_column → same as content but image is shown alongside text.
- Vary layouts naturally: do NOT repeat the same layout more than 2 times in a row. Do not default to content layout; make the visual rhythm feel intentionally designed.
- Aim for: 1 cover + 1-2 stats + 1 quote/timeline + 1 spotlight/card_grid + 1-3 content/two_column/image_focus + 1 closing.

CONTENT QUALITY:
- Every bullet must contain specific data: numbers, percentages, company names, dates, or named concepts.
- content field: write flowing analytical prose, not bullet-style fragments.
- title: short and impactful (noun phrase or question), under 40 characters.
- NEVER write "Lorem ipsum", "[insert]", or vague sentences like "다양한 요인이 있습니다".
{article_note}
{image_note}
{tone_rule}
- Language: ALL text fields in {lang_str}
"""


async def _call_llm_for_pages(
        prompt: str, articles: list, images: list[ImageMeta],
        page_count: int, page_count_auto: bool, language: str,
        reasoning: bool = False,
) -> dict:
    system = _build_system_prompt(language, page_count, page_count_auto, bool(articles), bool(images))

    # file:// 문서는 청크를 재조합하지 않고, 인덱싱 시 보관한 원문 전체를 조회한다.
    # 원문이 없는 이전 데이터만 호환성 차원에서 청크 재조합으로 보완한다.
    enriched_articles = []
    for a in articles:
        url = a.get("url", "")
        if url.startswith("file://") and not "::" in url:
            file_id = url.replace("file://", "")
            try:
                es = get_es()
                try:
                    original = await es.get(index=DOCUMENT_ORIGINALS_INDEX, id=file_id)
                    full_content = original.get("_source", {}).get("content", "")
                except Exception:
                    res = await es.search(
                        index=DOC_CHUNKS_INDEX,
                        body={"query": {"term": {"file_id": file_id}}, "sort": [{"chunk_index": "asc"}], "size": 50},
                        _source=["content", "chunk_index"]
                    )
                    chunks = [h["_source"].get("content", "") for h in res["hits"]["hits"]]
                    full_content = "\n\n".join(chunks) if chunks else a.get("content", "")
                    enriched_articles.append({**a, "content": full_content})
                else:
                    enriched_articles.append({**a, "content": full_content or a.get("content", "")})
                finally:
                    await es.close()
            except Exception:
                enriched_articles.append(a)
        else:
            enriched_articles.append(a)

    content_char_limit = PRESENTATION_CONTEXT_CHAR_LIMIT // max(len(enriched_articles), 1)
    context_docs = [
        {"title": a.get("title", ""), "content": a.get("content", "")[:content_char_limit],
         "source": a.get("source", ""), "url": a.get("url", ""), "score": 1.0}
        for a in enriched_articles
    ]

    img_desc = ""
    if images:
        img_desc = f"\n\n[첨부 이미지 {len(images)}장]\n"
        for img in images:
            img_desc += f"- Image {img.index}: {img.filename}\n"
        img_desc += "각 이미지를 분석하여 적절한 페이지에 배치하거나 내용에 반영하세요."

    # 이미지 임시 저장
    attachments = []
    for img in images:
        tmp_fn = f"pdf_tmp_{uuid.uuid4().hex[:8]}.{img.filename.split('.')[-1]}"
        tmp_path = FILES_DIR / tmp_fn
        tmp_path.write_bytes(base64.b64decode(img.data))
        attachments.append({"type": "image", "filename": tmp_fn, "_tmp": True})

    temperature_token = set_request_temperature_override(PRESENTATION_TEMPERATURE)
    try:
        raw, _ = await collect_llm_stream(
            question=f"{prompt}{img_desc}",
            context_docs=context_docs,
            system_prompt=system,
            attachments=attachments,
            timeout=300.0,
            format_instruction_override="",
            reasoning=reasoning,  # 프론트 추론 스위치 값 반영
            call_reason="presentation_generation",
        )
    finally:
        reset_request_temperature_override(temperature_token)
        for att in attachments:
            if att.get("_tmp"):
                try:
                    (FILES_DIR / att["filename"]).unlink(missing_ok=True)
                except Exception:
                    pass

    raw = raw.strip()

    # 마크다운 코드 펜스 제거 (```json ... ``` 또는 ``` ... ```)
    if raw.startswith("```"):
        lines = raw.splitlines()
        if lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        raw = "\n".join(lines).strip()

    # 직접 파싱 시도
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        # 에러 위치 전후 컨텍스트 로깅 (디버깅용)
        pos = e.pos
        snippet = raw[max(0, pos - 100):pos + 100]
        logger.error(
            "[pdf] JSON 파싱 실패: %s | pos=%d | 에러 주변:\n---\n%s\n---",
            e.msg, pos, snippet,
        )
        raise ValueError(f"LLM JSON 파싱 실패 — {e.msg} (line {e.lineno} col {e.colno})\n주변 텍스트: {snippet}")


# ── HTML 슬라이드 렌더러 ─────────────────────────────────────────────────────────
def _image_to_data_uri(img_meta: ImageMeta) -> str:
    mime = img_meta.type or "image/jpeg"
    return f"data:{mime};base64,{img_meta.data}"


def _render_bullets(bullets: list, accent2: str, body_color: str) -> str:
    """불릿 리스트 HTML 렌더링 (f-string 중첩 회피용 헬퍼)"""
    items = []
    for b in bullets:
        items.append(
            f'<div class="bullet-item" style="border-left:3px solid {accent2};">'
            f'<span style="color:{accent2};font-weight:700;">▸</span>{b}</div>'
        )
    return "\n  ".join(items)


def _render_page_html(page: dict, img_map: dict, palette: dict, total: int) -> str:
    layout   = page.get("layout", "content")
    title    = page.get("title", "")
    subtitle = page.get("subtitle") or ""
    content  = page.get("content") or ""
    bullets  = page.get("bullets") or []
    stats    = page.get("stats") or []
    quote    = page.get("quote") or ""
    img_idx  = page.get("image_index")
    img_pos  = page.get("image_position") or "right"
    caption  = page.get("image_caption") or ""
    page_num = page.get("index", 0) + 1
    has_img  = img_idx is not None and img_idx in img_map

    p       = palette
    accent  = p["accent"]
    accent2 = p["accent2"]
    img_uri = _image_to_data_uri(img_map[img_idx]) if has_img else ""

    # ── COVER ──
    if layout == "cover":
        tag_bg   = p.get("tag_bg", accent)
        tag_text = p.get("tag_text", "#fff")
        ct       = p.get("cover_text", "#fff")
        ct_sub   = p.get("cover_text_sub", "rgba(255,255,255,0.75)")
        side_bar = p.get("cover_side_bar", accent)
        deco_c   = p.get("cover_deco", accent2)
        year_c   = "rgba(0,0,0,0.25)" if ct == "#111111" else "rgba(255,255,255,0.3)"
        tech_tags = ""
        if bullets:
            tech_tags = "".join(
                f'<span class="cover-tag" style="background:{tag_bg};color:{tag_text};">{b}</span>'
                for b in bullets[:6]
            )
        return f"""
<div class="slide cover-slide" style="background:{p['cover_bg']};">
  <div class="cover-side-bar" style="background:{side_bar};"></div>
  <div class="cover-deco-circle" style="background:{deco_c};"></div>
  <div class="cover-inner">
    <h1 class="cover-title" style="color:{ct};">{title}</h1>
    {"<p class='cover-subtitle' style='color:" + ct_sub + ";'>" + subtitle + "</p>" if subtitle else ""}
    {"<p class='cover-desc' style='color:" + ct_sub + ";opacity:0.8;'>" + content + "</p>" if content else ""}
    {"<div class='cover-tags'>" + tech_tags + "</div>" if tech_tags else ""}
  </div>
  <div class="cover-year" style="color:{year_c};">{page_num} / {total}</div>
</div>"""

    # ── CLOSING ──
    if layout == "closing":
        ct      = p.get("cover_text", "#fff")
        ct_sub  = p.get("cover_text_sub", "rgba(255,255,255,0.75)")
        side_c  = p.get("cover_side_bar", accent2)
        deco_c  = p.get("cover_deco", accent)
        year_c  = "rgba(0,0,0,0.25)" if ct == "#111111" else "rgba(255,255,255,0.3)"
        # closing은 bullets/stats 무시 — 항상 심플 텍스트만
        closing_content = content or ("; ".join(bullets) if bullets else "")
        return f"""
<div class="slide cover-slide closing-slide" style="background:{p['cover_bg']};">
  <div class="cover-side-bar" style="background:{side_c};"></div>
  <div class="cover-deco-circle" style="background:{deco_c};"></div>
  <div class="cover-inner" style="text-align:center;align-items:center;">
    <div class="closing-icon" style="border-color:{accent};color:{accent};">✦</div>
    <h1 class="cover-title" style="color:{ct};font-size:2.4rem;">{title}</h1>
    {"<p class='cover-subtitle' style='color:" + ct_sub + ";'>" + subtitle + "</p>" if subtitle else ""}
    {"<p class='cover-desc' style='color:" + ct_sub + ";opacity:0.85;'>" + closing_content + "</p>" if closing_content else ""}
  </div>
  <div class="cover-year" style="color:{year_c};">{page_num} / {total}</div>
</div>"""

    # ── STATS ──
    if layout == "stats" and stats:
        cards = ""
        for s in stats[:4]:
            cards += f"""
<div class="stat-card">
  <div class="stat-accent-bar" style="background:{accent};"></div>
  <div class="stat-value" style="color:{accent};">{s.get("value","")}</div>
  <div class="stat-label">{s.get("label","")}</div>
  {"<div class='stat-desc'>" + s["desc"] + "</div>" if s.get("desc") else ""}
</div>"""
        return f"""
<div class="slide doc-slide">
  {_doc_header(page_num, total, title, accent, accent2)}
  <div class="doc-body">
    {"<p class='doc-intro'>" + content + "</p>" if content else ""}
    <div class="stats-row">{cards}</div>
  </div>
</div>"""

    # ── QUOTE ──
    if layout == "quote":
        return f"""
<div class="slide doc-slide">
  {_doc_header(page_num, total, title, accent, accent2)}
  <div class="doc-body">
    <div class="quote-wrap" style="border-left:4px solid {accent};">
      <div class="quote-mark" style="color:{accent};">"</div>
      <p class="quote-text">{quote}</p>
      {"<p class='quote-source'>" + content + "</p>" if content else ""}
    </div>
  </div>
</div>"""

    def _bullet_cards(card_class: str = "") -> str:
        return "".join(
            f'''<div class="{card_class} bullet-card">
  <span class="bullet-card-index" style="color:{accent};">{index + 1:02d}</span>
  <p>{bullet}</p>
</div>'''
            for index, bullet in enumerate(bullets)
        )

    # ── TIMELINE ──
    if layout == "timeline" and bullets:
        steps = "".join(
            f'''<div class="timeline-step">
  <div class="timeline-dot" style="background:{accent};"></div>
  <div class="timeline-content">
    <span class="timeline-number" style="color:{accent};">{index + 1:02d}</span>
    <p>{bullet}</p>
  </div>
</div>'''
            for index, bullet in enumerate(bullets)
        )
        return f"""
<div class="slide doc-slide timeline-slide">
  {_doc_header(page_num, total, title, accent, accent2)}
  <div class="doc-body">
    {"<p class='doc-intro timeline-intro'>" + content + "</p>" if content else ""}
    <div class="timeline-list">{steps}</div>
  </div>
</div>"""

    # ── SPOTLIGHT ──
    if layout == "spotlight":
        image_html = f'<img src="{img_uri}" class="spotlight-image" alt="{caption}"/>' if has_img else ""
        return f"""
<div class="slide doc-slide spotlight-slide">
  {_doc_header(page_num, total, title, accent, accent2)}
  <div class="doc-body spotlight-body">
    <div class="spotlight-main" style="border-color:{accent};">
      <span class="spotlight-label" style="color:{accent};">✦</span>
      <p>{content or subtitle}</p>
    </div>
    {image_html}
    {"<div class='spotlight-support'>" + _bullet_cards() + "</div>" if bullets else ""}
  </div>
</div>"""

    # ── CARD GRID ──
    if layout == "card_grid" and bullets:
        return f"""
<div class="slide doc-slide card-grid-slide">
  {_doc_header(page_num, total, title, accent, accent2)}
  <div class="doc-body">
    {"<p class='doc-intro'>" + content + "</p>" if content else ""}
    <div class="insight-grid">{_bullet_cards()}</div>
  </div>
</div>"""

    # ── CONTENT / TWO_COLUMN / IMAGE_FOCUS ──
    # accent 밝기에 따라 번호 텍스트 색상 결정 (앰버 계열은 다크 텍스트)
    _accent_hex = accent.lstrip('#')
    try:
        _r, _g, _b = int(_accent_hex[0:2],16), int(_accent_hex[2:4],16), int(_accent_hex[4:6],16)
        _luminance = (_r * 0.299 + _g * 0.587 + _b * 0.114)
        _num_text = "#1a1200" if _luminance > 150 else "#ffffff"
    except Exception:
        _num_text = "#ffffff"

    def _bullet_items(bl):
        rows = ""
        for i, b in enumerate(bl):
            rows += f'''<div class="bullet-row">
  <div class="bullet-num" style="background:{accent};color:{_num_text};">{i+1}</div>
  <div class="bullet-text">{b}</div>
</div>'''
        return rows

    if has_img and img_pos in ("right", "left"):
        text_html  = (f'<p class="doc-intro" style="margin-bottom:3mm;">{content}</p>' if content else "")
        text_html += f'<div class="bullet-list">{_bullet_items(bullets)}</div>' if bullets else ""
        img_html   = f'''<div class="img-panel">
  <img src="{img_uri}" class="doc-img" alt="{caption}"/>
  {"<p class='img-cap'>" + caption + "</p>" if caption else ""}
</div>'''
        if img_pos == "right":
            cols = f'<div class="two-col img-r"><div class="text-panel">{text_html}</div>{img_html}</div>'
        else:
            cols = f'<div class="two-col img-l">{img_html}<div class="text-panel">{text_html}</div></div>'
    else:
        body = ""
        if has_img and img_pos == "full":
            body += f'<div class="full-img-wrap"><img src="{img_uri}" class="doc-img-full" alt="{caption}"/></div>'
        if content:
            body += f'<p class="doc-intro">{content}</p>'
        if bullets:
            body += f'<div class="bullet-list">{_bullet_items(bullets)}</div>'
        if subtitle and not content and not bullets:
            body += f'<p class="doc-intro">{subtitle}</p>'
        cols = f'<div class="single-body">{body}</div>'

    return f"""
<div class="slide doc-slide">
  {_doc_header(page_num, total, title, accent, accent2)}
  <div class="doc-body">
    {cols}
  </div>
</div>"""


def _doc_header(page_num: int, total: int, title: str, accent: str, accent2: str) -> str:
    return f'''<div class="doc-header">
  <div class="doc-header-left">
    <div class="doc-page-badge" style="color:{accent};border-color:{accent};">{page_num:02d} / {total:02d}</div>
    <h2 class="doc-title">{title}</h2>
  </div>
  <div class="doc-header-line" style="background:linear-gradient(90deg,{accent} 0%,{accent2} 40%,transparent 100%);"></div>
</div>'''


def _ensure_visual_rhythm(pages: list[dict]) -> None:
    """LLM 출력이 단조로운 경우, 불릿 기반 페이지에 안전한 정보 레이아웃을 부여한다."""
    if len(pages) < 5:
        return

    existing_layouts = {page.get("layout") for page in pages}
    candidates = [
        page for page in pages[1:-1]
        if page.get("bullets") and page.get("layout") in {"content", "two_column"}
    ]
    if not candidates:
        return

    if "timeline" not in existing_layouts:
        candidates[len(candidates) // 2]["layout"] = "timeline"

    remaining = [page for page in candidates if page.get("layout") in {"content", "two_column"}]
    if "card_grid" not in existing_layouts and remaining:
        remaining[-1]["layout"] = "card_grid"


def _build_html(page_data: dict, images: list[ImageMeta], style: str) -> str:
    palette   = STYLE_PALETTES.get(style, STYLE_PALETTES["white"])
    img_map   = {img.index: img for img in images}
    pages     = page_data.get("pages", [])
    total     = len(pages)
    prs_title = page_data.get("presentation_title", "Presentation")

    # 첫 페이지 → cover, 마지막 페이지 → closing 강제
    if pages:
        pages[0]["layout"] = "cover"
        pages[-1]["layout"] = "closing"
        _ensure_visual_rhythm(pages)

    slides_html = "\n".join(_render_page_html(p, img_map, palette, total) for p in pages)

    p = palette
    # 다크 여부 판단
    slide_bg     = p.get("slide_bg", "#ffffff")
    title_color  = p.get("title_color", "#111111")
    body_color   = p.get("body_color", "#374151")
    muted_color  = p.get("muted_color", "#6b7280")
    card_bg      = p.get("card_bg", "#f9fafb")
    card_bg2     = p.get("card_bg2", card_bg)
    border_color = p.get("border_color", "#e5e7eb")
    accent       = p["accent"]
    accent2      = p["accent2"]

    return f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8"/>
<title>{prs_title}</title>
<link rel="preconnect" href="https://fonts.googleapis.com"/>
<link href="https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@300;400;500;600;700;800;900&display=swap" rel="stylesheet"/>
<style>
  :root {{
    --slide-bg:     {slide_bg};
    --title-color:  {title_color};
    --body-color:   {body_color};
    --muted-color:  {muted_color};
    --card-bg:      {card_bg};
    --card-bg2:     {card_bg2};
    --border-color: {border_color};
    --accent:       {accent};
    --accent2:      {accent2};
  }}
  *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    font-family: 'Noto Sans KR', 'Malgun Gothic', sans-serif;
    -webkit-print-color-adjust: exact;
    print-color-adjust: exact;
    background: var(--slide-bg);
    width: 210mm;
    margin: 0;
    padding: 0;
  }}
  .slide {{
    width: 210mm;
    min-height: 297mm;
    position: relative;
    page-break-after: always;
    page-break-inside: avoid;
    background: var(--slide-bg);
    overflow-x: hidden;
  }}
  /* ── COVER ── */
  .cover-slide {{
    display: flex;
    flex-direction: column;
    justify-content: flex-end;
    padding: 0 20mm 18%;
    position: relative;
    min-height: 297mm;
  }}
  .cover-side-bar {{
    position: absolute; left:0; top:0;
    width: 4mm; height: 100%;
  }}
  .cover-deco-circle {{
    position: absolute;
    width: 110mm; height: 110mm;
    border-radius: 50%;
    right: -20mm; top: -20mm;
    opacity: 0.4;
    pointer-events: none;
  }}
  .cover-inner {{
    position: relative; z-index: 2;
    display: flex; flex-direction: column;
    gap: 5mm; max-width: 168mm;
  }}
  .closing-slide {{
    justify-content: center;
    padding-bottom: 0;
  }}
  .closing-slide .cover-inner {{
    align-items: center; text-align: center; margin: 0 auto;
  }}

  .cover-title {{
    font-size: 3.4rem; font-weight: 900;
    line-height: 1.15;
    word-break: keep-all; overflow-wrap: break-word;
    text-shadow: 0 2px 20px rgba(0,0,0,0.2);
    max-width: 160mm;
  }}
  .cover-subtitle {{
    font-size: 1rem; font-weight: 400;
    line-height: 1.65; word-break: keep-all;
  }}
  .cover-desc {{
    font-size: 0.83rem; line-height: 1.8;
    word-break: keep-all; max-width: 200mm;
  }}
  .cover-tags {{
    display: flex; flex-wrap: wrap; gap: 2mm; margin-top: 2mm;
  }}
  .cover-tag {{
    font-size: 0.68rem; font-weight: 700;
    padding: 1.5mm 4mm; border-radius: 20mm;
    letter-spacing: 0.03em;
  }}
  .cover-year {{
    position: absolute; bottom: 5mm; right: 8mm;
    font-size: 0.7rem; font-weight: 500;
  }}
  .closing-icon {{
    font-size: 1.5rem; width: 13mm; height: 13mm;
    border-radius: 50%; border: 2px solid;
    display: flex; align-items: center; justify-content: center;
    margin-bottom: 5mm;
  }}
  /* ── DOC SLIDE ── */
  .doc-slide {{
    display: flex; flex-direction: column;
    background: var(--slide-bg); min-height: 297mm;
    padding: 0 0 10mm 0;
  }}
  .doc-header {{
    flex-shrink: 0;
    padding: 12mm 16mm 0;
  }}
  .doc-header-left {{
    display: flex; align-items: center; gap: 4mm; margin-bottom: 2mm;
  }}
  .doc-page-badge {{
    font-size: 0.7rem; font-weight: 700; letter-spacing: 0.12em;
    border: 1.5px solid; border-radius: 20mm;
    padding: 1.5mm 4mm; white-space: nowrap; flex-shrink: 0;
  }}
  .doc-title {{
    font-size: 1.7rem; font-weight: 800; color: var(--title-color);
    line-height: 1.25; word-break: keep-all; overflow-wrap: break-word;
  }}
  .doc-header-line {{
    height: 2px; width: 100%; border-radius: 1mm;
    margin-top: 3mm;
  }}
  .doc-body {{
    flex: 1;
    padding: 8mm 16mm 12mm;
    display: flex; flex-direction: column; gap: 5mm;
  }}
  .doc-intro {{
    font-size: 0.95rem; line-height: 1.9;
    color: var(--body-color); word-break: keep-all;
    overflow-wrap: break-word; flex-shrink: 0;
  }}
  /* ── 번호 불릿 ── */
  .bullet-list {{
    display: flex; flex-direction: column; gap: 2.5mm; flex: 1;
  }}
  .bullet-row {{
    display: flex; align-items: center;
    gap: 4mm; padding: 4.5mm 5mm;
    background: var(--card-bg); border-radius: 2.5mm;
    border: 1px solid var(--border-color);
  }}
  .bullet-num {{
    width: 7mm; height: 7mm; border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    font-size: 0.7rem; font-weight: 800; flex-shrink: 0; line-height: 1;
  }}
  .bullet-text {{
    font-size: 0.93rem; line-height: 1.6; color: var(--body-color);
    word-break: keep-all; overflow-wrap: break-word; flex: 1;
  }}
  /* ── STATS ── */
  .stats-row {{
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(35mm, 1fr));
    gap: 4mm;
  }}
  .stat-card {{
    background: var(--card-bg); border: 1px solid var(--border-color);
    border-radius: 3mm; padding: 7mm 4mm 6mm;
    text-align: center; position: relative; overflow: hidden;
    min-width: 0;
  }}
  .stat-accent-bar {{
    position: absolute; top:0; left:0; right:0; height: 3px;
  }}
  .stat-value {{
    font-size: clamp(1rem, 3.5vw, 1.7rem); font-weight: 900;
    line-height: 1.2; margin-bottom: 3mm; word-break: keep-all;
    overflow-wrap: break-word; hyphens: auto;
  }}
  .stat-label {{
    font-size: 0.77rem; font-weight: 700; color: var(--title-color); margin-bottom: 1.5mm;
  }}
  .stat-desc {{
    font-size: 0.7rem; color: var(--muted-color); line-height: 1.5; word-break: keep-all;
  }}
  /* ── QUOTE ── */
  .quote-wrap {{
    flex: 1; padding: 5mm 6mm; border-radius: 2mm;
    background: var(--card-bg2);
    display: flex; flex-direction: column; justify-content: center; gap: 3mm;
  }}
  .quote-mark {{
    font-size: 2.5rem; font-family: Georgia, serif;
    font-weight: 900; line-height: 0.7;
  }}
  .quote-text {{
    font-size: 1.1rem; font-weight: 600;
    line-height: 1.8; color: var(--title-color); font-style: italic; word-break: keep-all;
  }}
  .quote-source {{
    font-size: 0.8rem; color: var(--body-color); line-height: 1.7; word-break: keep-all;
  }}
  /* ── TIMELINE ── */
  .timeline-intro {{ max-width: 155mm; }}
  .timeline-list {{ position:relative; display:flex; flex-direction:column; gap:0; margin:2mm 0 0 4mm; padding:0 0 0 9mm; }}
  .timeline-list::before {{ content:''; position:absolute; left:2.7mm; top:4mm; bottom:4mm; width:1px; background:var(--border-color); }}
  .timeline-step {{ position:relative; padding:0 0 5mm; min-height:19mm; }}
  .timeline-step:last-child {{ padding-bottom:0; }}
  .timeline-dot {{ position:absolute; left:-7.4mm; top:3.2mm; width:5mm; height:5mm; border:2px solid var(--slide-bg); border-radius:50%; box-shadow:0 0 0 1px var(--border-color); }}
  .timeline-content {{ display:flex; align-items:flex-start; gap:4mm; padding:3.5mm 5mm; background:var(--card-bg); border:1px solid var(--border-color); border-radius:2.5mm; }}
  .timeline-number {{ min-width:8mm; font-size:.74rem; font-weight:900; letter-spacing:.08em; }}
  .timeline-content p {{ font-size:.91rem; line-height:1.6; color:var(--body-color); word-break:keep-all; }}
  /* ── SPOTLIGHT ── */
  .spotlight-body {{ justify-content:center; gap:6mm; }}
  .spotlight-main {{ border-left:4px solid; padding:5mm 7mm 5mm 8mm; background:var(--card-bg2); border-radius:0 3mm 3mm 0; }}
  .spotlight-label {{ display:block; font-size:1.1rem; line-height:1; margin-bottom:3mm; }}
  .spotlight-main p {{ font-size:1.35rem; font-weight:750; line-height:1.55; color:var(--title-color); word-break:keep-all; }}
  .spotlight-image {{ width:100%; max-height:65mm; object-fit:cover; border-radius:3mm; }}
  .spotlight-support {{ display:grid; grid-template-columns:repeat(2, minmax(0, 1fr)); gap:3mm; }}
  /* ── CARD GRID ── */
  .insight-grid {{ display:grid; grid-template-columns:repeat(2, minmax(0, 1fr)); gap:4mm; align-content:start; }}
  .bullet-card {{ min-height:31mm; padding:5mm; background:var(--card-bg); border:1px solid var(--border-color); border-radius:3mm; }}
  .bullet-card-index {{ display:block; font-size:.72rem; font-weight:900; letter-spacing:.09em; margin-bottom:4mm; }}
  .bullet-card p {{ font-size:.9rem; line-height:1.65; color:var(--body-color); word-break:keep-all; }}
  .spotlight-support .bullet-card {{ min-height:0; padding:3.5mm 4mm; }}
  .spotlight-support .bullet-card-index {{ display:none; }}
  /* ── TWO-COL ── */
  .two-col {{ display: flex; gap: 5mm; flex: 1; overflow: hidden; }}
  .text-panel {{ flex: 1; display: flex; flex-direction: column; gap: 2.5mm; overflow: hidden; }}
  .img-panel {{ flex: 1; display: flex; flex-direction: column; align-items: center; justify-content: center; gap: 2mm; }}
  .doc-img {{ max-width:100%; max-height:110mm; object-fit:contain; border-radius:2mm; box-shadow:0 3px 16px rgba(0,0,0,0.1); }}
  .full-img-wrap {{ width:100%; display:flex; justify-content:center; margin-bottom:3mm; flex-shrink:0; }}
  .doc-img-full {{ max-width:100%; max-height:65mm; object-fit:contain; border-radius:2mm; }}
  .img-cap {{ font-size:0.7rem; color:#888; text-align:center; font-style:italic; }}
  .single-body {{ display:flex; flex-direction:column; gap:3mm; flex:1; overflow:hidden; }}
  @media print {{
    body {{ background:#fff; }}
    .slide {{ page-break-after:always; page-break-inside:avoid; }}
  }}
  @page {{ size:A4 portrait; margin:0; }}
</style>
</head>
<body>
{slides_html}
</body>
</html>"""


async def _html_to_pdf(html_path: Path, output_path: Path):
    """Playwright로 HTML → PDF 변환"""
    from playwright.async_api import async_playwright

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(args=["--no-sandbox", "--disable-setuid-sandbox"])
        page = await browser.new_page()
        await page.goto(f"file://{html_path}", wait_until="networkidle", timeout=30000)
        # 폰트 로딩 대기
        await page.wait_for_timeout(1500)
        await page.pdf(
            path=str(output_path),
            format="A4",
            print_background=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
        )
        await browser.close()


async def _html_to_pptx(html_path: Path, output_path: Path) -> None:
    """Render each HTML slide to a PNG and package the rendered slides as a PPTX.

    Rendering preserves the same layouts, typography, and theme as the PDF output.
    """
    from playwright.async_api import async_playwright

    with tempfile.TemporaryDirectory(prefix="vyact_pptx_") as temp_dir:
        image_paths: list[Path] = []
        async with async_playwright() as pw:
            browser = await pw.chromium.launch(args=["--no-sandbox", "--disable-setuid-sandbox"])
            page = await browser.new_page(viewport={"width": 1240, "height": 1754}, device_scale_factor=1)
            await page.goto(f"file://{html_path}", wait_until="networkidle", timeout=30000)
            await page.wait_for_timeout(1000)
            slides = page.locator(".slide")
            for index in range(await slides.count()):
                image_path = Path(temp_dir) / f"slide_{index + 1}.png"
                await slides.nth(index).screenshot(path=str(image_path))
                image_paths.append(image_path)
            await browser.close()

        presentation = Presentation()
        presentation.slide_width = Inches(8.27)
        presentation.slide_height = Inches(11.69)
        blank_layout = presentation.slide_layouts[6]
        for image_path in image_paths:
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        await asyncio.to_thread(presentation.save, str(output_path))


def _sse_step(step: int, total: int, message_key: str, pct: int, **message_params: int | str) -> str:
    return sse(json.dumps({
        "step": step,
        "total": total,
        "message_key": message_key,
        "message_params": message_params,
    }), "progress", pct)


@router.post("/pdf/generate")
async def generate_pdf(req: PdfGenerateRequest):
    STEPS = 7

    async def event_stream():
        html_path = None
        n_articles = len([a for a in req.articles if not a.get("url","").startswith(("memo://","file://"))])
        n_imgs = len(req.images)
        page_label = "AI 자동" if req.page_count_auto else f"{req.page_count}페이지"
        logger.info("[pdf] 생성 시작 — 기사:%d 이미지:%d 페이지설정:%s", n_articles, n_imgs, page_label)
        user_ts = datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")

        try:
            # ── Step 1: 소스 준비 ──
            yield _sse_step(1, STEPS, "sourcePreparation", 5, articleCount=n_articles, imageCount=n_imgs)
            logger.info("[pdf] step1 소스 준비")

            # ── Step 2: 문서 컨텍스트 구성 ──
            yield _sse_step(2, STEPS, "contextPreparation", 12)

            # ── Step 3: AI 스토리·콘텐츠 생성 ──
            yield _sse_step(3, STEPS, "slideDesignAuto" if req.page_count_auto else "slideDesign", 20, pageCount=req.page_count)
            logger.info("[pdf] step3 LLM 호출 시작")

            try:
                page_data = await _call_llm_for_pages(
                    prompt=req.prompt, articles=req.articles, images=req.images,
                    page_count=req.page_count, page_count_auto=req.page_count_auto,
                    language=req.language,
                    reasoning=req.reasoning,
                )
            except Exception as llm_err:
                logger.error("[pdf] step3 LLM/JSON 파싱 실패: %s", llm_err, exc_info=True)
                yield sse(json.dumps({"error": f"슬라이드 구조 생성 실패: {llm_err}"}), "error")
                return

            n_pages = len(page_data.get("pages", []))
            logger.info("[pdf] step3 LLM 완료 — %d페이지 확정", n_pages)

            # ── Step 4: 슬라이드 구성 확인 ──
            yield _sse_step(4, STEPS, "layoutCheck", 55, slideCount=n_pages)

            # ── Step 5: HTML 렌더링 ──
            yield _sse_step(5, STEPS, "htmlRendering", 64, pageCount=n_pages)
            logger.info("[pdf] step5 HTML 렌더링 시작")

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            uid = str(uuid.uuid4())[:6]

            html_content = _build_html(page_data, req.images, req.style)
            html_path = FILES_DIR / f"pdf_tmp_{uid}.html"
            html_path.write_text(html_content, encoding="utf-8")
            logger.info("[pdf] step5 HTML 완료 (%d bytes)", len(html_content))

            # ── Step 6: 선택한 출력 형식으로 렌더링 ──
            output_format = "pptx" if req.output_format == "pptx" else "pdf"
            format_label = "PPTX" if output_format == "pptx" else "PDF"
            yield _sse_step(6, STEPS, "outputRendering", 76, outputFormat=output_format, pageCount=n_pages)
            logger.info("[pdf] step6 %s 변환 시작", format_label)

            temp_output_path = FILES_DIR / f"pdf_tmp_{uid}.{output_format}"
            if output_format == "pptx":
                await _html_to_pptx(html_path, temp_output_path)
            else:
                await _html_to_pdf(html_path, temp_output_path)
            logger.info("[pdf] step6 %s 변환 완료", format_label)

            # ── Step 7: 저장 및 마무리 ──
            yield _sse_step(7, STEPS, "saving", 92)
            logger.info("[pdf] step7 저장 시작")

            # 이미지 영구 저장 (재편집용)
            saved_image_filenames = []
            for img in req.images:
                ext = img.filename.rsplit(".", 1)[-1] if "." in img.filename else "jpg"
                img_fn = f"pdf_{uid}_{img.index}.{ext}"
                (IMAGES_DIR / img_fn).write_bytes(base64.b64decode(img.data))
                saved_image_filenames.append(img_fn)

            conv_id = req.conv_id or str(uuid.uuid4())
            prs_title = page_data.get("presentation_title", req.prompt[:30])

            # 제목 기반 파일명 (특수문자 제거, 공백→언더스코어, 최대 40자)
            import re as _re
            _bad_chars = r'[\/*?:<>|]'
            safe_title = _re.sub(_bad_chars, '', prs_title)
            safe_title = _re.sub(r'\s+', '_', safe_title.strip())[:40].rstrip('_')
            filename = f"{safe_title}_{timestamp}.{output_format}"
            output_path = FILES_DIR / filename
            temp_output_path.rename(output_path)

            style_name = STYLE_PALETTES.get(req.style, {}).get("name", req.style)

            # 소스 타입별 카운트
            n_articles = sum(1 for a in req.articles if not a.get("url","").startswith(("memo://","file://")))
            n_memos    = sum(1 for a in req.articles if a.get("url","").startswith("memo://"))
            n_docs     = sum(1 for a in req.articles if a.get("url","").startswith("file://"))
            source_parts = []
            if n_articles: source_parts.append(f"기사 {n_articles}개")
            if n_memos:    source_parts.append(f"메모 {n_memos}개")
            if n_docs:     source_parts.append(f"문서 {n_docs}개")
            source_str = " · ".join(source_parts) if source_parts else "소스 없음"

            answer_text = (
                    f"📄 **{prs_title}** {format_label}가 생성되었습니다.\n\n"
                    f"{n_pages}페이지 · 스타일: {style_name} · "
                    f"{source_str} · 이미지 {len(req.images)}장 활용"
                    + (" · 자동 페이지 수 선택" if req.page_count_auto else "")
            )

            pdf_params = {
                "prompt": req.prompt,
                "page_count": req.page_count,
                "page_count_auto": req.page_count_auto,
                "language": req.language,
                "style": req.style,
                "output_format": output_format,
                "articles": [
                    {
                        "url": a.get("url", ""),
                        "title": a.get("title", ""),
                        "source": a.get("source", ""),
                        "indexed_at": a.get("indexed_at", ""),
                        "file_id": a.get("file_id"),
                        # 인덱스에 저장하지 않는 프레젠테이션 첨부는 재편집 시에도
                        # 동일한 원문 컨텍스트를 사용할 수 있도록 함께 보관한다.
                        "content": a.get("content", "") if a.get("url", "").startswith("attachment://") else "",
                    }
                    for a in req.articles
                ],
                "image_filenames": saved_image_filenames,
            }

            now = datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")
            await save_conversation(conv_id, req.messages + [
                {"role": "user", "content": f"/presentation {req.prompt}", "timestamp": user_ts},
                {
                    "role": "assistant",
                    "content": answer_text,
                    "timestamp": now,
                    "model": await get_model_name(),
                    "pdf_file": filename,
                    "pdf_params": pdf_params,
                },
            ])

            yield sse(json.dumps({
                "answer": answer_text,
                "filename": filename,
                "conv_id": conv_id,
                "pdf_params": pdf_params,
            }), "done", 100)
            logger.info("[pdf] 생성 완료 — %s (%d페이지)", filename, n_pages)

        except (asyncio.CancelledError, GeneratorExit):
            logger.info("[pdf] 클라이언트 연결 종료 — PDF 생성 중단")
        except Exception as e:
            import traceback
            tb = traceback.format_exc()
            logger.error("[pdf] 생성 실패: %s\n%s", e, tb)
            yield sse(json.dumps({"error": str(e)}), "error")
        finally:
            if html_path and html_path.exists():
                html_path.unlink(missing_ok=True)

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@router.get("/pdf/download/{filename}")
async def download_pdf(filename: str):
    safe_name = Path(filename).name
    filepath = FILES_DIR / safe_name
    if not filepath.exists():
        from fastapi import HTTPException
        raise HTTPException(404, "파일을 찾을 수 없습니다")
    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation" if safe_name.endswith(".pptx") else "application/pdf"
    return FileResponse(
        str(filepath),
        media_type=media_type,
        filename=safe_name,
    )
